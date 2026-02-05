"""
Generate PowerPoint presentation summarizing the Hypertension Microsimulation Model.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
ACCENT_GREEN = RGBColor(0, 128, 0)
ACCENT_RED = RGBColor(192, 0, 0)
GRAY = RGBColor(89, 89, 89)


def add_title_slide(title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add dark blue header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full blue background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # Title centered
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, bullet_points, two_column=False):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if two_column and len(bullet_points) > 4:
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True

        mid = len(bullet_points) // 2
        for i, (text, level) in enumerate(bullet_points[:mid]):
            if i == 0:
                p = left_tf.paragraphs[0]
            else:
                p = left_tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18) if level == 0 else Pt(16)
            p.font.color.rgb = GRAY if level > 0 else DARK_BLUE
            p.level = level
            p.space_before = Pt(6)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True

        for i, (text, level) in enumerate(bullet_points[mid:]):
            if i == 0:
                p = right_tf.paragraphs[0]
            else:
                p = right_tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18) if level == 0 else Pt(16)
            p.font.color.rgb = GRAY if level > 0 else DARK_BLUE
            p.level = level
            p.space_before = Pt(6)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, (text, level) in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20) if level == 0 else Pt(18)
            p.font.color.rgb = GRAY if level > 0 else DARK_BLUE
            p.font.bold = (level == 0)
            p.level = level
            p.space_before = Pt(8)

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = Inches(12.3)
    table_height = Inches(min(5.5, 0.5 * num_rows))

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5),
        table_width, table_height
    ).table

    # Set column widths
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

            # Highlight dominated cells
            if "DOMINATED" in str(value):
                p.font.color.rgb = ACCENT_RED
                p.font.bold = True

    return slide


# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
add_title_slide(
    "Hypertension Cardiorenal Microsimulation Model",
    "Cost-Effectiveness Analysis of IXA-001 vs Spironolactone | Version 4.0"
)

# =============================================================================
# SLIDE 2: Executive Summary
# =============================================================================
add_content_slide(
    "Executive Summary",
    [
        ("Model Purpose: Evaluate IXA-001 (novel aldosterone synthase inhibitor) in resistant hypertension", 0),
        ("Individual-level microsimulation with dual cardiac-renal disease branches", 1),
        ("Key Innovation (v4.0): Atrial fibrillation as aldosterone-specific outcome", 0),
        ("PA patients have 12× baseline AF risk; IXA-001 provides 60% AF risk reduction", 1),
        ("Primary Finding: PA patients are the optimal target population", 0),
        ("ICER: $245,441/QALY (6.7% price reduction achieves $150K threshold)", 1),
        ("Essential HTN is contraindicated (dominated: higher cost, worse outcomes)", 1),
        ("Economic Perspective: Societal (direct + indirect costs)", 0),
        ("Event cost savings offset 72-83% of drug premium", 1),
    ]
)

# =============================================================================
# SECTION: Background
# =============================================================================
add_section_slide("Background & Target Population")

# =============================================================================
# SLIDE 3: Target Population
# =============================================================================
add_content_slide(
    "Target Population: Resistant Hypertension",
    [
        ("Definition: BP ≥130/80 mmHg despite ≥3 antihypertensives (including diuretic)", 0),
        ("At maximally tolerated doses with confirmed adherence", 1),
        ("Prevalence: 10-15% of treated hypertensive patients", 0),
        ("~11,000 per 1 million health plan members", 1),
        ("Why Microsimulation Required (vs. Markov cohort models):", 0),
        ("25-35% have prior CV events (individual history tracking essential)", 1),
        ("30-40% have CKD (dual cardiac-renal pathways)", 1),
        ("40-50% have diabetes (accelerated progression)", 1),
        ("15-20% have Primary Aldosteronism (IXA-001 target)", 1),
    ]
)

# =============================================================================
# SLIDE 4: Primary Aldosteronism - The IXA-001 Target
# =============================================================================
add_content_slide(
    "Primary Aldosteronism: The IXA-001 Target",
    [
        ("15-20% of resistant HTN patients have Primary Aldosteronism (PA)", 0),
        ("Autonomous aldosterone production drives hypertension", 1),
        ("Pathophysiology driving elevated risk:", 0),
        ("Cardiac fibrosis → 2.05× heart failure risk (Monticone 2018)", 1),
        ("Atrial remodeling → 12× atrial fibrillation risk", 1),
        ("Renal fibrosis → 1.80× ESRD risk", 1),
        ("Vascular stiffness → 1.50× stroke risk", 1),
        ("Why IXA-001 works better in PA:", 0),
        ("Aldosterone synthase inhibition directly addresses pathophysiology", 1),
        ("70% enhanced BP response (modifier: 1.70×)", 1),
        ("60% AF risk reduction (vs 40% with spironolactone)", 1),
    ]
)

# =============================================================================
# SECTION: Model Architecture
# =============================================================================
add_section_slide("Model Architecture")

# =============================================================================
# SLIDE 5: Model Overview
# =============================================================================
add_content_slide(
    "Model Structure Overview",
    [
        ("Model Type: Individual-Level State-Transition Microsimulation (IL-STM)", 0),
        ("Cycle Length: Monthly (captures acute events, rapid renal transitions)", 0),
        ("Time Horizon: Lifetime (up to age 100) or configurable (e.g., 20 years)", 0),
        ("Perspective: Healthcare system OR Societal (configurable)", 0),
        ("Discount Rate: 3% per annum (costs and QALYs)", 0),
        ("Comparators:", 0),
        ("IXA-001 (novel ASI): $500/month", 1),
        ("Spironolactone (generic MRA): $15/month", 1),
        ("Population: Resistant HTN with stratification by secondary etiology", 0),
    ]
)

# =============================================================================
# SLIDE 6: Dual Disease Branch Architecture
# =============================================================================
add_content_slide(
    "Dual Disease Branch Architecture",
    [
        ("Cardiac Branch (simultaneous with renal):", 0),
        ("No Acute Event → MI / Stroke / HF / AF → Chronic states → CV Death", 1),
        ("Recurrent events tracked individually", 1),
        ("Renal Branch (simultaneous with cardiac):", 0),
        ("CKD Stage 1-2 → 3a → 3b → 4 → ESRD → Renal Death", 1),
        ("Continuous eGFR decline (not discrete jumps)", 1),
        ("Dynamic Risk Factors (updated monthly):", 0),
        ("SBP: Age drift + stochastic variation - treatment effect", 1),
        ("eGFR: Baseline decline × diabetes × BP effect", 1),
        ("Competing Risks: CV death vs Renal death vs Other-cause mortality", 0),
    ]
)

# =============================================================================
# SLIDE 7: Key Model Features (v4.0)
# =============================================================================
add_content_slide(
    "Key Model Features (Version 4.0)",
    [
        ("1. Individual patient tracking (history, adherence, time since events)", 0),
        ("2. Dynamic SBP model: SBP(t+1) = SBP(t) + β_age + ε - treatment_effect", 0),
        ("3. Enhanced eGFR decline: age-stratified × diabetes × BP interaction", 0),
        ("4. Four-dimensional risk stratification:", 0),
        ("GCUA (age ≥60) / EOCRI (age 18-59) phenotypes", 1),
        ("KDIGO renal risk matrix + Framingham CVD risk", 1),
        ("5. NEW: Atrial fibrillation as aldosterone-specific outcome", 0),
        ("12× baseline risk for PA; differential treatment effects", 1),
        ("6. NEW: Societal perspective with indirect costs", 0),
        ("Productivity loss, absenteeism, chronic disability", 1),
        ("7. Safety monitoring: Hyperkalemia management for MRA patients", 0),
    ]
)

# =============================================================================
# SECTION: Risk Calculation
# =============================================================================
add_section_slide("Risk Calculation Methodology")

# =============================================================================
# SLIDE 8: Event Probability Calculation
# =============================================================================
add_content_slide(
    "How Event Probabilities Are Calculated",
    [
        ("Layer 1: AHA PREVENT Equations (base 10-year CVD risk)", 0),
        ("Inputs: Age, sex, SBP, eGFR, diabetes, smoking, cholesterol, BMI", 1),
        ("Layer 2: Conversion to Monthly Probability", 0),
        ("p_annual = 1 - (1 - p_10yr)^0.1 → p_monthly = 1 - (1 - p_annual)^(1/12)", 1),
        ("Layer 3: Prior Event Multipliers", 0),
        ("Prior MI: 2.5× | Prior Stroke: 3.0× | Prior HF: 2.0×", 1),
        ("Layer 4: Phenotype-Specific Modifiers (PA, OSA, RAS, Pheo)", 0),
        ("PA example: MI 1.40×, Stroke 1.50×, HF 2.05×, AF 3.0×", 1),
        ("Layer 5: Treatment Risk Reduction", 0),
        ("BP reduction → risk reduction via efficacy coefficients", 1),
    ]
)

# =============================================================================
# SLIDE 9: Phenotype Risk Modifiers Table
# =============================================================================
add_table_slide(
    "Secondary HTN Etiology Risk Modifiers",
    ["Etiology", "MI", "Stroke", "HF", "ESRD", "AF", "Death"],
    [
        ("Primary Aldosteronism", "1.40×", "1.50×", "2.05×", "1.80×", "3.0×", "1.60×"),
        ("Obstructive Sleep Apnea", "1.25×", "1.35×", "1.40×", "1.30×", "1.5×", "1.35×"),
        ("Renal Artery Stenosis", "1.30×", "1.40×", "1.35×", "1.60×", "1.2×", "1.40×"),
        ("Pheochromocytoma", "1.50×", "1.60×", "1.30×", "1.10×", "1.3×", "1.50×"),
        ("Essential HTN", "1.0×", "1.0×", "1.0×", "1.0×", "1.0×", "1.0×"),
    ],
    col_widths=[2.5, 1.5, 1.5, 1.5, 1.5, 1.5, 1.5]
)

# =============================================================================
# SLIDE 10: Treatment Response Modifiers
# =============================================================================
add_table_slide(
    "Treatment Response Modifiers by Etiology",
    ["Etiology", "IXA-001 Response", "Spironolactone Response", "Clinical Rationale"],
    [
        ("Primary Aldosteronism", "1.70×", "1.40×", "Direct aldosterone pathway targeting"),
        ("Obstructive Sleep Apnea", "1.20×", "1.15×", "Aldosterone elevation in OSA"),
        ("Renal Artery Stenosis", "0.90×", "0.85×", "Volume/RAAS-driven HTN"),
        ("Pheochromocytoma", "0.70×", "0.65×", "Catecholamine-driven, not aldosterone"),
        ("Essential HTN", "1.0×", "1.0×", "Baseline response"),
    ],
    col_widths=[2.5, 2.2, 2.5, 4.5]
)

# =============================================================================
# SECTION: Results
# =============================================================================
add_section_slide("Cost-Effectiveness Results")

# =============================================================================
# SLIDE 11: Study Design
# =============================================================================
add_content_slide(
    "Study Design",
    [
        ("Population: 2,000 resistant hypertension patients", 0),
        ("Time Horizon: 20 years", 0),
        ("Perspective: Societal (direct + indirect costs)", 0),
        ("Discount Rate: 3% per annum (costs and QALYs)", 0),
        ("Comparators:", 0),
        ("IXA-001: $500/month (novel aldosterone synthase inhibitor)", 1),
        ("Spironolactone: $15/month (generic MRA)", 1),
        ("Willingness-to-Pay Threshold: $150,000/QALY", 0),
        ("Subgroups stratified by secondary HTN etiology:", 0),
        ("PA (21%), OSA (15%), RAS (11%), Essential (52%)", 1),
    ]
)

# =============================================================================
# SLIDE 12: Main Results Table
# =============================================================================
add_table_slide(
    "Cost-Effectiveness Results by Subgroup (20-Year)",
    ["Subgroup", "N", "Δ Cost", "Δ QALYs", "ICER ($/QALY)"],
    [
        ("Primary Aldosteronism", "425", "+$20,550", "+0.084", "$245,441"),
        ("Obstructive Sleep Apnea", "305", "+$33,245", "+0.129", "$258,370"),
        ("Renal Artery Stenosis", "221", "+$25,906", "+0.092", "$281,298"),
        ("Essential HTN", "1,030", "+$28,568", "-0.062", "DOMINATED"),
    ],
    col_widths=[3.5, 1.5, 2.5, 2.0, 2.8]
)

# =============================================================================
# SLIDE 13: Event Prevention (PA Subgroup)
# =============================================================================
add_table_slide(
    "Event Prevention Analysis: PA Subgroup (n=425)",
    ["Event", "IXA-001", "Spironolactone", "Events Prevented", "Rate Ratio"],
    [
        ("Myocardial Infarction", "21", "39", "+18", "0.54"),
        ("Stroke", "27", "48", "+21", "0.56"),
        ("Heart Failure", "24", "41", "+17", "0.59"),
        ("New Atrial Fibrillation", "225", "258", "+33", "0.87"),
        ("CV Deaths", "271", "270", "-1", "1.00"),
    ],
    col_widths=[3.0, 2.0, 2.5, 2.5, 2.0]
)

# =============================================================================
# SLIDE 14: AF Prevention (New Outcome)
# =============================================================================
add_table_slide(
    "Atrial Fibrillation Prevention (NEW v4.0 Outcome)",
    ["Subgroup", "AF (IXA-001)", "AF (Spiro)", "AF Prevented", "Clinical Significance"],
    [
        ("Primary Aldosteronism", "225", "258", "+33", "Primary aldosterone target"),
        ("Obstructive Sleep Apnea", "72", "77", "+5", "Moderate benefit"),
        ("Renal Artery Stenosis", "41", "44", "+3", "Minimal benefit"),
        ("Essential HTN", "238", "238", "0", "No differential"),
    ],
    col_widths=[3.0, 2.0, 2.0, 2.0, 3.3]
)

# =============================================================================
# SLIDE 15: Cost Component Analysis
# =============================================================================
add_content_slide(
    "Cost Component Analysis",
    [
        ("20-Year Drug Cost Premium: $116,400", 0),
        ("($500 - $15) × 12 months × 20 years", 1),
        ("Event Cost Savings Partially Offset Premium:", 0),
        ("PA: ~$95,850 savings (82% offset) → Net: $20,550", 1),
        ("OSA: ~$83,155 savings (71% offset) → Net: $33,245", 1),
        ("RAS: ~$90,494 savings (78% offset) → Net: $25,906", 1),
        ("Essential: ~$87,832 savings (75% offset) → Net: $28,568", 1),
        ("Key Finding: Event prevention offsets 72-83% of drug cost premium", 0),
        ("PA has highest offset due to more events prevented", 1),
    ]
)

# =============================================================================
# SLIDE 16: Threshold Pricing
# =============================================================================
add_table_slide(
    "Threshold Price Analysis (at $150,000/QALY WTP)",
    ["Subgroup", "Current ICER", "Threshold Price", "Price Reduction"],
    [
        ("Primary Aldosteronism", "$245,441", "$467/month", "6.7%"),
        ("Obstructive Sleep Apnea", "$258,370", "$442/month", "11.6%"),
        ("Renal Artery Stenosis", "$281,298", "$450/month", "10.1%"),
        ("Essential HTN", "DOMINATED", "N/A", "N/A"),
    ],
    col_widths=[3.5, 2.5, 2.5, 2.5]
)

# =============================================================================
# SLIDE 17: Key Findings
# =============================================================================
add_content_slide(
    "Key Findings & Clinical Recommendations",
    [
        ("1. PA Patients Are the Primary Value Driver", 0),
        ("Largest event reduction: 18 MI, 21 stroke, 17 HF, 33 AF prevented", 1),
        ("Smallest price cut needed (6.7%) to achieve CE threshold", 1),
        ("2. OSA Patients Show Meaningful Benefit", 0),
        ("Highest QALY gain (+0.129) due to enhanced treatment response", 1),
        ("3. Essential HTN is Contraindicated for IXA-001", 0),
        ("DOMINATED: Higher cost AND worse outcomes (negative QALYs)", 1),
        ("No aldosterone-specific pathophysiology to target", 1),
        ("4. Pheochromocytoma Requires Different Treatment", 0),
        ("ASI ineffective; requires alpha/beta blockade", 1),
    ]
)

# =============================================================================
# SECTION: Economic Parameters
# =============================================================================
add_section_slide("Economic Parameters")

# =============================================================================
# SLIDE 18: Costs
# =============================================================================
add_table_slide(
    "Cost Parameters",
    ["Category", "Item", "Cost", "Source"],
    [
        ("Event (Acute)", "MI", "$25,000", "Hospitalization + cath lab"),
        ("Event (Acute)", "Ischemic Stroke", "$15,200", "ICU + imaging + rehab"),
        ("Event (Acute)", "Hemorrhagic Stroke", "$22,500", "Neurosurgery + ICU"),
        ("Event (Acute)", "Heart Failure", "$18,000", "Inpatient stay"),
        ("Event (Acute)", "New AF", "$8,500", "Cardioversion + anticoag"),
        ("State (Monthly)", "ESRD", "$7,500", "Dialysis"),
        ("State (Monthly)", "Chronic HF", "$1,500", "HF clinic + devices"),
        ("State (Monthly)", "Chronic AF", "$708", "DOAC + monitoring"),
        ("Drug", "IXA-001", "$500/mo", "Novel ASI"),
        ("Drug", "Spironolactone", "$15/mo", "Generic MRA"),
    ],
    col_widths=[2.0, 3.0, 2.0, 5.3]
)

# =============================================================================
# SLIDE 19: Utilities
# =============================================================================
add_table_slide(
    "Utility Parameters (EQ-5D Scale)",
    ["Category", "Condition", "Value", "Reference"],
    [
        ("Baseline", "Age 40-50", "0.87", "Sullivan 2006"),
        ("Baseline", "Age 60-70", "0.81", "Sullivan 2006"),
        ("Baseline", "Age 80+", "0.72", "Sullivan 2006"),
        ("Chronic Disutility", "Post-MI", "-0.12", "Lacey 2003"),
        ("Chronic Disutility", "Post-Stroke", "-0.18", "Luengo-Fernandez 2013"),
        ("Chronic Disutility", "Chronic HF", "-0.15", "Calvert 2021"),
        ("Chronic Disutility", "Chronic AF", "-0.05", "Dorian 2000"),
        ("Chronic Disutility", "ESRD", "-0.35", "Wasserfallen 2004"),
        ("Acute Disutility", "New AF (1 month)", "-0.15", "Dorian 2000"),
    ],
    col_widths=[2.5, 3.0, 1.5, 5.3]
)

# =============================================================================
# SLIDE 20: Indirect Costs (Societal)
# =============================================================================
add_table_slide(
    "Indirect Costs (Societal Perspective)",
    ["Component", "Value", "Application"],
    [
        ("Daily Wage (US)", "$240", "Working age <65 only"),
        ("MI Absenteeism", "7 days", "One-time acute event"),
        ("Stroke Absenteeism", "30 days", "One-time acute event"),
        ("HF Absenteeism", "5 days", "One-time acute event"),
        ("Stroke Chronic Disability", "20% annual wage", "Ongoing productivity loss"),
        ("HF Chronic Disability", "15% annual wage", "Ongoing productivity loss"),
    ],
    col_widths=[3.5, 2.5, 6.3]
)

# =============================================================================
# SECTION: Conclusions
# =============================================================================
add_section_slide("Conclusions & Recommendations")

# =============================================================================
# SLIDE 21: Conclusions
# =============================================================================
add_content_slide(
    "Conclusions",
    [
        ("1. IXA-001 is not cost-effective at $500/month for any subgroup", 0),
        ("At $150,000/QALY WTP threshold", 1),
        ("2. PA patients are the optimal target population", 0),
        ("Highest event prevention, smallest price reduction needed (6.7%)", 1),
        ("3. AF prevention is a key value differentiator", 0),
        ("33 AF events prevented in PA subgroup alone", 1),
        ("4. Essential HTN is a contraindication", 0),
        ("Dominated outcomes: higher cost, worse QALYs", 1),
        ("5. Tiered pricing strategy recommended", 0),
        ("~$467/month for PA, ~$445/month for OSA/RAS", 1),
        ("6. Event cost savings offset 72-83% of drug premium", 0),
        ("Net budget impact more favorable than gross drug cost", 1),
    ]
)

# =============================================================================
# SLIDE 22: Recommended Pricing Strategy
# =============================================================================
add_table_slide(
    "Recommended Tiered Pricing Strategy",
    ["Tier", "Subgroups", "Recommended Price", "Expected ICER", "Rationale"],
    [
        ("Tier 1", "Primary Aldosteronism", "$467/month", "~$150,000", "Core target population"),
        ("Tier 2", "OSA, RAS", "$445/month", "~$150,000", "Secondary responders"),
        ("Exclude", "Essential, Pheo", "N/A", "N/A", "No benefit / contraindicated"),
    ],
    col_widths=[1.5, 3.0, 2.5, 2.0, 3.3]
)

# =============================================================================
# SLIDE 23: Sensitivity Analysis Drivers
# =============================================================================
add_table_slide(
    "Sensitivity Analysis: Key ICER Drivers",
    ["Parameter", "Base Case", "Impact on ICER"],
    [
        ("IXA-001 monthly cost", "$500", "±$100 → ICER ±$95,000"),
        ("IXA-001 SBP reduction", "24 mmHg", "±5 mmHg → ICER ±$60,000"),
        ("PA treatment modifier", "1.70×", "±0.2 → ICER ±$40,000"),
        ("AF event cost", "$8,500", "±$3,000 → ICER ±$15,000"),
        ("Time horizon", "20 years", "40 years → ICER -30%"),
        ("Discount rate", "3%", "5% → ICER +20%"),
    ],
    col_widths=[4.0, 3.0, 5.3]
)

# =============================================================================
# SLIDE 24: Next Steps
# =============================================================================
add_content_slide(
    "Next Steps",
    [
        ("Model Validation:", 0),
        ("External validation against PA cohort studies", 1),
        ("PSA with 10,000 iterations for robust confidence intervals", 1),
        ("Regulatory/HTA Preparation:", 0),
        ("CHEERS 2022 compliant technical report", 1),
        ("Budget impact analysis for US and UK payers", 1),
        ("Clinical Development:", 0),
        ("Confirm PA diagnostic pathway for patient selection", 1),
        ("Real-world evidence collection on AF outcomes", 1),
        ("Pricing & Access:", 0),
        ("Negotiate tiered pricing with demonstration of PA diagnosis", 1),
        ("Outcomes-based contracts tied to BP response", 1),
    ]
)

# =============================================================================
# SLIDE 25: References
# =============================================================================
add_content_slide(
    "Key References",
    [
        ("Clinical:", 0),
        ("Monticone S, et al. JACC 2018 - PA cardiovascular risk", 1),
        ("Khan SS, et al. Circulation 2024 - PREVENT equations", 1),
        ("Laffin LJ, et al. NEJM 2023 - Baxdrostat efficacy", 1),
        ("Williams B, et al. Lancet 2015 - PATHWAY-2 trial", 1),
        ("Health Economics:", 0),
        ("Briggs A, et al. Oxford 2006 - Decision modelling", 1),
        ("Sanders GD, et al. JAMA 2016 - CEA recommendations", 1),
        ("Husereau D, et al. BMJ 2022 - CHEERS 2022", 1),
        ("AF-Specific:", 0),
        ("Dorian P, et al. JACC 2000 - AF quality of life", 1),
        ("Kim MH, et al. Circ CV Qual Outcomes 2011 - AF costs", 1),
    ],
    two_column=True
)

# =============================================================================
# Final slide
# =============================================================================
add_title_slide(
    "Thank You",
    "Hypertension Microsimulation Model v4.0 | Questions?"
)

# Save presentation
output_path = "/home/dan/Genesis Interview/hypertension_microsim/IXA-001_CEA_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
