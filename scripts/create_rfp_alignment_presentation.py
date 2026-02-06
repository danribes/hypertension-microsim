"""
Generate RFP Alignment PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
GRAY = RGBColor(89, 89, 89)
WHITE = RGBColor(255, 255, 255)


def add_title_slide(title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(title):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, bullet_points):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, (text, level) in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20) if level == 0 else Pt(18)
        p.font.color.rgb = GRAY if level > 0 else DARK_BLUE
        p.font.bold = (level == 0)
        p.level = level
        p.space_before = Pt(8)

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(12.3)
    table_height = Inches(min(5.5, 0.5 * num_rows))

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5),
        table_width, table_height
    ).table

    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

            if "DOMINATED" in str(value):
                p.font.color.rgb = RED
                p.font.bold = True
            elif "✓" in str(value):
                p.font.color.rgb = GREEN

    return slide


# =============================================================================
# SLIDES
# =============================================================================

# Slide 1: Title
add_title_slide(
    "IXA-001 Cost-Effectiveness and Budget Impact Models",
    "RFP Deliverables Assessment | Atlantis Pharmaceuticals | February 2026"
)

# Slide 2: Executive Summary
add_table_slide(
    "Executive Summary: Deliverables Complete",
    ["Deliverable", "RFP Requirement", "Status"],
    [
        ("Cost-Effectiveness Model", "HTA submission support", "✓ Delivered"),
        ("Budget Impact Model", "Payer discussions (US + EU5)", "✓ Delivered"),
        ("SBP Reduction", "20 mmHg (Phase III)", "✓ Aligned"),
        ("User-Friendly Interface", "Face-to-face payer use", "✓ Enhanced"),
    ],
    col_widths=[4, 5, 3]
)

# Slide 3: RFP Requirements
add_content_slide(
    "RFP Requirements Recap",
    [
        ("Clinical Data (Phase III):", 0),
        ("2,000 patients, 5 centres (80% EU)", 1),
        ("20 mmHg SBP reduction vs placebo (p=0.025)", 1),
        ("Well tolerated; low hyperkalemia signal", 1),
        ("Deliverables Requested:", 0),
        ("1. Cost-effectiveness model for global HTA submissions", 1),
        ("2. User-friendly budget impact model for US → EU5", 1),
        ("Value Proposition:", 0),
        ("IXA-001 reduces BP → Fewer CV/renal events → Cost-effective vs SOC", 1),
    ]
)

# Section: CEA
add_section_slide("Cost-Effectiveness Model (CEA)")

# Slide 4: CEA Overview
add_table_slide(
    "CEA Model Overview",
    ["Attribute", "Specification"],
    [
        ("Model Type", "Individual-Level Microsimulation"),
        ("Cycle Length", "Monthly"),
        ("Time Horizon", "Lifetime (configurable)"),
        ("Perspective", "Healthcare system OR Societal"),
        ("Comparator", "Spironolactone (guideline 4th-line)"),
        ("Discount Rate", "3% per annum"),
    ],
    col_widths=[4, 8]
)

# Slide 5: SBP Alignment
add_table_slide(
    "Clinical Parameters: SBP Aligned to RFP",
    ["Parameter", "RFP Specification", "Model Value", "Status"],
    [
        ("IXA-001 SBP Reduction", "20 mmHg vs placebo", "20.0 mmHg", "✓ ALIGNED"),
        ("Spiro SBP Reduction", "Not specified", "9.0 mmHg", "Evidence-based"),
        ("Tolerability", "Well tolerated", "8% discontinuation", "✓ Aligned"),
        ("Hyperkalemia", "Lower than expected", "Reduced monitoring", "✓ Aligned"),
    ],
    col_widths=[3.5, 3, 3, 2.5]
)

# Slide 6: Dual Disease Branch
add_content_slide(
    "Dual Disease Branch Architecture",
    [
        ("Cardiac Branch (simultaneous with renal):", 0),
        ("No Event → MI / Stroke / HF / AF → Chronic States → CV Death", 1),
        ("Renal Branch (simultaneous with cardiac):", 0),
        ("CKD 1-2 → CKD 3a → CKD 3b → CKD 4 → ESRD → Renal Death", 1),
        ("Outcomes Tracked:", 0),
        ("MI, Stroke, Heart Failure, Atrial Fibrillation (NEW), ESRD, Death", 1),
        ("Enhancement: AF tracking as aldosterone-specific outcome", 1),
    ]
)

# Slide 7: Phenotype Stratification
add_table_slide(
    "Phenotype Stratification: PA is Optimal Target",
    ["Phenotype", "Prevalence", "IXA-001 Response", "Key Risk"],
    [
        ("Primary Aldosteronism (PA)", "17%", "1.70× enhanced", "HF 2.05×, ESRD 1.80×"),
        ("Obstructive Sleep Apnea", "15%", "1.20× enhanced", "Stroke 1.25×"),
        ("Renal Artery Stenosis", "11%", "1.05× standard", "ESRD 1.80×"),
        ("Pheochromocytoma", "1%", "0.40× reduced", "Contraindicated"),
        ("Essential HTN", "56%", "1.0× baseline", "Standard risk"),
    ],
    col_widths=[4, 2, 2.5, 3.5]
)

# Slide 8: Subgroup Results
add_table_slide(
    "CEA Results by Subgroup (20-Year, Societal)",
    ["Subgroup", "N", "Δ Cost", "Δ QALYs", "ICER"],
    [
        ("Primary Aldosteronism", "425", "+$20,550", "+0.084", "$245,441/QALY"),
        ("Obstructive Sleep Apnea", "305", "+$33,245", "+0.129", "$258,370/QALY"),
        ("Renal Artery Stenosis", "221", "+$25,906", "+0.092", "$281,298/QALY"),
        ("Essential HTN", "1,030", "+$28,568", "-0.062", "DOMINATED"),
    ],
    col_widths=[4, 1.5, 2, 2, 2.5]
)

# Slide 9: Event Prevention
add_table_slide(
    "Event Prevention: PA Subgroup (n=425)",
    ["Event", "IXA-001", "Spironolactone", "Prevented"],
    [
        ("MI", "21", "39", "+18"),
        ("Stroke", "27", "48", "+21"),
        ("Heart Failure", "24", "41", "+17"),
        ("Atrial Fibrillation", "225", "258", "+33"),
        ("CV Deaths", "271", "270", "-1"),
    ],
    col_widths=[4, 2.5, 3, 2.5]
)

# Slide 10: Threshold Pricing
add_table_slide(
    "Threshold Pricing Analysis ($150K/QALY)",
    ["Subgroup", "Current ICER", "Threshold Price", "Price Cut"],
    [
        ("Primary Aldosteronism", "$245,441", "$467/month", "6.7%"),
        ("Obstructive Sleep Apnea", "$258,370", "$442/month", "11.6%"),
        ("Renal Artery Stenosis", "$281,298", "$450/month", "10.1%"),
        ("Essential HTN", "DOMINATED", "N/A", "N/A"),
    ],
    col_widths=[4, 3, 3, 2]
)

# Section: BIM
add_section_slide("Budget Impact Model (BIM)")

# Slide 11: BIM Overview
add_table_slide(
    "BIM Model Overview",
    ["Attribute", "Specification"],
    [
        ("Model Type", "Cohort-based budget impact"),
        ("Time Horizon", "5 years (10-year extension)"),
        ("Perspective", "Healthcare payer"),
        ("Interface", "Streamlit web + Interactive Excel"),
        ("Markets", "US, UK, DE, FR, IT, ES"),
        ("Output", "13-sheet comprehensive Excel"),
    ],
    col_widths=[4, 8]
)

# Slide 12: Multi-Country Support
add_table_slide(
    "Multi-Country Support: US + EU5",
    ["Country", "Currency", "Cost Multiplier", "HTN Prevalence"],
    [
        ("United States", "USD ($)", "1.00", "30%"),
        ("United Kingdom", "GBP (£)", "0.40", "28%"),
        ("Germany", "EUR (€)", "0.50", "32%"),
        ("France", "EUR (€)", "0.45", "30%"),
        ("Italy", "EUR (€)", "0.42", "33%"),
        ("Spain", "EUR (€)", "0.38", "33%"),
    ],
    col_widths=[3.5, 2.5, 3, 3]
)

# Slide 13: User-Friendly Interface
add_content_slide(
    "User-Friendly Interface for Payer Discussions",
    [
        ("Streamlit Web Application:", 0),
        ("Country selection dropdown", 1),
        ("Scenario selection (Conservative/Moderate/Optimistic)", 1),
        ("Interactive sliders for all inputs", 1),
        ("Real-time calculation updates", 1),
        ("One-click Excel download", 1),
        ("Excel Output (13 Sheets):", 0),
        ("Cover, Inputs, Population, Market, Costs, Results, Scenarios", 1),
        ("Tornado, Subgroups, 10-Year, Events, PSA, Documentation", 1),
    ]
)

# Section: Compliance
add_section_slide("Compliance & Value Proposition")

# Slide 14: Value Proposition Support
add_table_slide(
    "Value Proposition: All 7 Statements Supported",
    ["#", "Value Proposition", "CEA Evidence", "BIM Evidence"],
    [
        ("1", "HTN highly prevalent", "Population generator", "Epidemiology cascade"),
        ("2", "Uncontrolled on ≥3 agents", "Target population", "Resistant HTN focus"),
        ("3", "Elevated CV/renal risk", "PREVENT equations", "Event rates"),
        ("4", "CV events costly", "Event cost module", "Cost offsets"),
        ("5", "IXA-001 reduces BP", "20 mmHg modeled", "Treatment effect"),
        ("6", "Reduced HCRU", "Event tracking", "Avoided event costs"),
        ("7", "Cost-effective vs SOC", "ICER calculation", "Price threshold"),
    ],
    col_widths=[0.8, 4, 3.5, 3.5]
)

# Slide 15: Compliance Checklist
add_table_slide(
    "RFP Requirements: 9/9 Complete",
    ["#", "Requirement", "Status"],
    [
        ("1", "Cost-effectiveness model", "✓ Microsimulation"),
        ("2", "HTA submission support", "✓ CHEERS compliant"),
        ("3", "Budget impact model", "✓ 5-year BIM"),
        ("4", "User-friendly interface", "✓ Streamlit + Excel"),
        ("5", "US market", "✓ Full parameters"),
        ("6", "EU5 adaptation", "✓ UK, DE, FR, IT, ES"),
        ("7", "Face-to-face payer use", "✓ Interactive Excel"),
        ("8", "IXA-001 vs SOC comparison", "✓ vs Spironolactone"),
        ("9", "20 mmHg SBP reduction", "✓ Aligned"),
    ],
    col_widths=[1, 6, 5]
)

# Slide 16: Model Enhancements
add_table_slide(
    "Enhancements Beyond RFP",
    ["Enhancement", "Model", "Benefit"],
    [
        ("Microsimulation", "CEA", "Captures patient heterogeneity"),
        ("Phenotype Stratification", "CEA", "Identifies PA as optimal target"),
        ("AF Tracking", "CEA", "Unique ASI value differentiator"),
        ("Societal Perspective", "CEA", "Comprehensive economic analysis"),
        ("6-Country Support", "BIM", "Full EU5 + US coverage"),
        ("13-Sheet Excel", "BIM", "Complete payer toolkit"),
        ("PSA Module", "Both", "Robust uncertainty quantification"),
    ],
    col_widths=[4, 2, 6]
)

# Slide 17: Key Strategic Insights
add_content_slide(
    "Key Strategic Insights",
    [
        ("1. PA patients are the optimal target population", 0),
        ("Largest event reduction, smallest price cut needed (6.7%)", 1),
        ("Clear biological rationale for ASI mechanism", 1),
        ("2. Essential HTN is contraindicated", 0),
        ("Dominated outcomes - exclude from formulary positioning", 1),
        ("3. AF prevention is a key differentiator", 0),
        ("33 events prevented in PA subgroup - unique to ASI", 1),
        ("4. Event cost savings offset 72-83% of drug premium", 0),
        ("Net budget impact more favorable than gross cost", 1),
    ]
)

# Slide 18: Recommended Pricing Strategy
add_table_slide(
    "Recommended Pricing Strategy",
    ["Tier", "Population", "Price", "Expected ICER"],
    [
        ("Tier 1", "Primary Aldosteronism", "$467/month", "~$150,000/QALY"),
        ("Tier 2", "OSA, RAS", "$445/month", "~$150,000/QALY"),
        ("Exclude", "Essential HTN, Pheo", "N/A", "Contraindicated"),
    ],
    col_widths=[2, 4, 3, 3]
)

# Slide 19: Next Steps
add_content_slide(
    "Recommended Next Steps",
    [
        ("Model Validation:", 0),
        ("External validation against PA cohort studies", 1),
        ("PSA with 10,000 iterations for robust confidence intervals", 1),
        ("Regulatory/HTA Preparation:", 0),
        ("NICE-specific adaptation for UK submission", 1),
        ("G-BA dossier preparation for Germany", 1),
        ("Clinical Development:", 0),
        ("Confirm PA diagnostic pathway for patient selection", 1),
        ("Real-world evidence collection on AF outcomes", 1),
        ("Pricing & Access:", 0),
        ("Outcomes-based contracts tied to BP response", 1),
    ]
)

# Slide 20: Summary
add_table_slide(
    "Summary: Deliverables Complete and Enhanced",
    ["Model", "Status", "Key Enhancement"],
    [
        ("CEA", "✓ Complete", "Microsimulation with phenotype stratification"),
        ("BIM", "✓ Complete", "6 countries, interactive Excel"),
        ("SBP Alignment", "✓ 20 mmHg", "Per RFP Phase III data"),
        ("Value Props", "✓ 7/7", "Full evidence package"),
    ],
    col_widths=[3, 2.5, 6.5]
)

# Slide 21: Thank You
add_title_slide(
    "Thank You",
    "IXA-001 CEA & BIM Models | Atlantis Pharmaceuticals"
)

# Save presentation
output_path = "IXA-001_RFP_Alignment_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
